# -*- coding: utf-8 -*-
"""紹介資料を PowerPoint の形にする（dist/センゴク盤_紹介資料.pptx）

   中身は src/data/shiryo.js ひとつから取る（build/shiryo.json 経由）。

   紙と Word は読ませるための形だが、こちらは映すための形である。同じ字を
   そのまま貼れば、字が小さすぎて後ろの席から読めない。ゆえにこう分けた。

     章の扉　　  章ごとに一枚。題と副題だけを大きく置く
     文の枚　　  一節につき一枚。長い文は箇条に割り、収まらなければ次の枚へ送る
     表の枚　　  表はそのまま一枚。行が多ければ割る
     数の枚　　  数字を大きく並べる
     絵の枚　　  写しは全面に置き、説きを下に添える

   一枚に載せる行数を決め打ちにしてある。溢れたら次の枚へ送るので、字が
   潰れることはない。
"""
import json, os, re
from pptx import Presentation
from pptx.util import Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
DATA = json.load(open(os.path.join(ROOT, 'build', 'shiryo.json'), encoding='utf-8'))
DIST = os.path.join(ROOT, 'dist')
絵蔵 = os.path.join(DIST, 'tebiki', '大')
図蔵 = os.path.join(DIST, 'tebiki', 'zu')

墨 = RGBColor(0x1C, 0x1E, 0x22)
紙 = RGBColor(0xF4, 0xF1, 0xE8)
薄 = RGBColor(0x6E, 0x6A, 0x62)
灰 = RGBColor(0x8A, 0x84, 0x78)
本 = RGBColor(0x33, 0x31, 0x2E)
藍 = RGBColor(0x2F, 0x5D, 0x8C)
明朝 = "Hiragino Mincho ProN"
ゴシック = "Hiragino Sans"

prs = Presentation()
prs.slide_width, prs.slide_height = Cm(33.867), Cm(19.05)   # 16:9
W, H = prs.slide_width, prs.slide_height
空 = prs.slide_layouts[6]

余 = Cm(2.2)
幅 = W - 余 * 2

def 枚(地=紙):
    s = prs.slides.add_slide(空)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = 地
    return s

def 字(s, 文, x, y, w, h, 大=18, 色=本, 書=明朝, 太=False, 揃=PP_ALIGN.LEFT,
       行間=1.35, 間=None, 中央=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if 中央: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = 揃
    p.line_spacing = 行間
    r = p.add_run(); r.text = 文
    f = r.font; f.size = Pt(大); f.color.rgb = 色; f.name = 書; f.bold = 太
    if 間 is not None:
        from pptx.oxml.ns import qn
        r._r.get_or_add_rPr().set('spc', str(int(間 * 100)))
    return tb

def 線(s, x, y, w, 太=Pt(1.6), 色=墨):
    from pptx.enum.shapes import MSO_SHAPE
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, 太)
    ln.fill.solid(); ln.fill.fore_color.rgb = 色; ln.line.fill.background()
    ln.shadow.inherit = False
    return ln

def 見出し(s, 題, 副=None):
    字(s, 題, 余, Cm(1.3), 幅, Cm(1.4), 大=27, 色=墨, 書=明朝)
    if 副:
        字(s, 副, 余, Cm(2.65), 幅, Cm(0.8), 大=12, 色=灰, 書=ゴシック, 間=1.5)
    線(s, 余, Cm(3.55), 幅)

# ------------------------------------------------------------------ 表紙
s = 枚(墨)
題PNG = os.path.join(ROOT, 'build', 'logo-title.png')
if os.path.exists(題PNG):
    w = Cm(17)
    s.shapes.add_picture(題PNG, int((W - w) / 2), Cm(4.6), width=w)
字(s, DATA['資料名'], 0, Cm(11.4), W, Cm(1), 大=17, 色=紙, 揃=PP_ALIGN.CENTER, 間=4)
字(s, DATA['添え書き'], 0, Cm(12.9), W, Cm(0.9), 大=11, 色=RGBColor(0xA8, 0xA2, 0x96),
   揃=PP_ALIGN.CENTER, 間=1.5)

# ------------------------------------------------------------------ 目次
s = 枚()
見出し(s, '目次')
半 = len(DATA['資料']) - len(DATA['資料']) // 2
for i, c in enumerate(DATA['資料']):
    列 = 0 if i < 半 else 1
    行 = i if i < 半 else i - 半
    x = 余 + 列 * (幅 // 2)
    y = Cm(4.6) + Cm(1.55) * 行
    字(s, '%d' % (i + 1), x, y + Cm(0.12), Cm(1.2), Cm(0.8), 大=12, 色=灰, 書=ゴシック)
    字(s, c['題'], x + Cm(1.2), y, Cm(9), Cm(0.9), 大=17, 色=墨)
    if c.get('副'):
        字(s, c['副'], x + Cm(1.2), y + Cm(0.92), Cm(12), Cm(0.7), 大=10, 色=灰, 書=ゴシック)

# ------------------------------------------------------------------ 本文
def 章の扉(i, c):
    s = 枚(墨)
    字(s, '%02d' % (i + 1), 余, Cm(6.2), Cm(4), Cm(1.6), 大=13, 色=RGBColor(0x6E, 0x6A, 0x62), 書=ゴシック, 間=3)
    字(s, c['題'], 余, Cm(7.5), 幅, Cm(2.4), 大=42, 色=紙)
    if c.get('副'):
        字(s, c['副'], 余, Cm(10.5), 幅, Cm(1), 大=14, 色=RGBColor(0xA8, 0xA2, 0x96), 書=ゴシック, 間=2.5)
    線(s, 余, Cm(12.2), Cm(4), Pt(2), 紙)

import unicodedata

def 字幅(t):
    """和文は一字、欧文は半字として数える。行数の見積もりに使う。"""
    return sum(1.0 if unicodedata.east_asian_width(c) in 'WFA' else 0.5 for c in t)

def 行数(t, 箱幅Emu, 大きさ):
    """その大きさで、その幅の箱に何行で入るか。和文は字高＝字幅とみなす。"""
    一行 = max(1.0, Emu(箱幅Emu).pt / 大きさ)
    return max(1, int(-(-字幅(t) // 一行)))

def 割る(文, 幅字=46):
    """長い一文を、読める長さに折る。句点で切り、それでも長ければ読点で切る。"""
    出 = []
    for 文2 in re.split(r'(?<=。)', 文):
        文2 = 文2.strip()
        if not 文2: continue
        if len(文2) <= 幅字 * 2:
            出.append(文2); continue
        買 = ''
        for 片 in re.split(r'(?<=、)', 文2):
            if len(買) + len(片) > 幅字 * 2 and 買:
                出.append(買); 買 = 片
            else:
                買 += 片
        if 買: 出.append(買)
    return 出

def 箇条の枚(題, 副, 行, 一枚=6, 大=15):
    """箇条を、決めた数ずつ枚に分けて置く。溢れたら次の枚へ送る。"""
    for k in range(0, len(行), 一枚):
        s = 枚()
        見出し(s, 題, 副)
        y = Cm(4.7)
        w = 幅 - Cm(0.9)
        for t in 行[k:k + 一枚]:
            n = 行数(t, w, 大)
            h = Emu(int(Pt(大 * 1.4).emu * n))          # 行間ぶんを掛けた高さ
            字(s, '—', 余, y + Cm(0.06), Cm(0.8), Cm(0.7), 大=13, 色=灰, 書=ゴシック)
            字(s, t, 余 + Cm(0.9), y, w, h, 大=大, 色=本, 行間=1.4)
            y += h + Cm(0.5)
        if len(行) > 一枚:
            字(s, '%d／%d' % (k // 一枚 + 1, (len(行) + 一枚 - 1) // 一枚),
               余, H - Cm(1.5), 幅, Cm(0.6), 大=9, 色=灰, 書=ゴシック, 揃=PP_ALIGN.RIGHT)

def 表の枚(題, 副, t, 一枚=6):
    行 = t.get('行') or []
    if not 行: return
    列 = len(行[0])
    for k in range(0, len(行), 一枚):
        塊 = 行[k:k + 一枚]
        s = 枚()
        見出し(s, 題, 副)
        y0 = Cm(4.5)
        頭 = t.get('頭') or []
        # 欄の幅。左の見出し欄は狭く、残りを等分する
        左 = Cm(6.4) if 列 == 2 else Cm(4.6)
        残 = (幅 - 左) // max(1, 列 - 1)
        xs = [余] + [余 + 左 + 残 * j for j in range(列 - 1)]
        ws = [左] + [残] * (列 - 1)
        if any(頭):
            for j, h in enumerate(頭):
                if h: 字(s, h, xs[j], y0, ws[j] - Cm(0.4), Cm(0.6), 大=10, 色=薄, 書=ゴシック, 間=1.5)
            線(s, 余, y0 + Cm(0.75), 幅, Pt(1.2))
            y0 += Cm(1.15)
        tbl = prs.slides[-1]
        y = y0
        余白 = (H - Cm(2.2) - y0) / max(1, len(塊))
        for r in 塊:
            for j, v in enumerate(r):
                大きさ = 15 if j == 0 else 12.5
                色 = 墨 if j == 0 else 本
                字(s, v, xs[j], y, ws[j] - Cm(0.4), 余白 - Cm(0.2),
                   大=大きさ, 色=色, 太=(j == 0), 行間=1.35)
            線(s, 余, y + 余白 - Cm(0.32), 幅, Pt(0.5), RGBColor(0xD8, 0xD2, 0xC4))
            y += 余白
        if len(行) > 一枚:
            字(s, '%d／%d' % (k // 一枚 + 1, (len(行) + 一枚 - 1) // 一枚),
               余, H - Cm(1.5), 幅, Cm(0.6), 大=9, 色=灰, 書=ゴシック, 揃=PP_ALIGN.RIGHT)

def 数の枚(題, 副, k):
    for b in range(0, len(k), 6):
        塊 = k[b:b + 6]
        s = 枚()
        見出し(s, 題, 副)
        for i, 項 in enumerate(塊):
            名, 値, 添 = (項 + ['', '', ''])[:3]
            列, 行 = i % 3, i // 3
            w = 幅 // 3
            x = 余 + 列 * w
            y = Cm(5.2) + 行 * Cm(5.4)
            字(s, 名, x, y, w - Cm(0.6), Cm(0.6), 大=11, 色=灰, 書=ゴシック, 間=1.5)
            # 長い値は字を落とす。「2,240 〜 167,802 石」のような一行が潰れぬように。
            大き = 34
            while 大き > 15 and 行数(値, w - Cm(0.6), 大き) > 1: 大き -= 2
            字(s, 値, x, y + Cm(0.75), w - Cm(0.6), Cm(1.6), 大=大き, 色=墨)
            if 添: 字(s, 添, x, y + Cm(2.7), w - Cm(0.6), Cm(1.8), 大=10.5, 色=薄, 行間=1.45)

def 絵の枚(題, 副, key, 説):
    f = os.path.join(絵蔵, key + '.jpg')
    if not os.path.exists(f): return
    s = 枚()
    見出し(s, 題, 副)
    h = H - Cm(6.0)
    s.shapes.add_picture(f, 0, 0, height=h)
    pic = s.shapes[-1]
    pic.left = int((W - pic.width) / 2); pic.top = Cm(4.4)
    if pic.width > 幅:
        pic.height = int(pic.height * 幅 / pic.width); pic.width = 幅
        pic.left = 余
    if 説: 字(s, 説, 余, H - Cm(1.9), 幅, Cm(1), 大=11, 色=薄, 揃=PP_ALIGN.CENTER)

def 図の枚(題, 副, z, 一枚=4):
    for k in range(0, len(z), 一枚):
        塊 = z[k:k + 一枚]
        s = 枚()
        見出し(s, 題, 副)
        y = Cm(4.6)
        高 = (H - Cm(6.2)) / 一枚
        for key, 名, 説 in 塊:
            f = os.path.join(図蔵, key + '.jpg')
            if os.path.exists(f):
                s.shapes.add_picture(f, 余, y, height=min(高 - Cm(0.4), Cm(2.6)))
            字(s, 名, 余 + Cm(4.4), y, Cm(5), Cm(0.8), 大=16, 色=墨, 太=True)
            字(s, 説, 余 + Cm(4.4), y + Cm(0.95), 幅 - Cm(4.4), 高 - Cm(1.1),
               大=11.5, 色=本, 行間=1.4)
            y += 高
        if len(z) > 一枚:
            字(s, '%d／%d' % (k // 一枚 + 1, (len(z) + 一枚 - 1) // 一枚),
               余, H - Cm(1.5), 幅, Cm(0.6), 大=9, 色=灰, 書=ゴシック, 揃=PP_ALIGN.RIGHT)

for i, c in enumerate(DATA['資料']):
    章の扉(i, c)
    題 = c['題']
    for s2 in c['節']:
        副 = s2.get('見出し') or c.get('副')
        if s2.get('絵'):
            絵の枚(題, 副, s2['絵'], s2.get('絵の説'))
        文 = []
        for t in s2.get('文') or []:
            文 += 割る(t)
        if 文: 箇条の枚(題, 副, 文, 一枚=5)
        if s2.get('表'): 表の枚(題, 副, s2['表'])
        if s2.get('数'): 数の枚(題, 副, s2['数'])
        if s2.get('図'): 図の枚(題, 副, s2['図'])
        if s2.get('箇条'): 箇条の枚(題, 副, s2['箇条'], 一枚=6)

# ------------------------------------------------------------------ 結び
s = 枚(墨)
印PNG = os.path.join(ROOT, 'build', 'logo-seal.png')
if os.path.exists(印PNG):
    w = Cm(5.4)
    s.shapes.add_picture(印PNG, int((W - w) / 2), Cm(5.4), width=w)
字(s, DATA['題名'], 0, Cm(11.4), W, Cm(1.4), 大=30, 色=紙, 揃=PP_ALIGN.CENTER, 間=4)
字(s, DATA['副題'], 0, Cm(13.3), W, Cm(0.8), 大=11, 色=RGBColor(0xA8, 0xA2, 0x96),
   揃=PP_ALIGN.CENTER, 間=6)

out = os.path.join(DIST, '%s_%s.pptx' % (DATA['題名'], DATA['資料名']))
prs.save(out)
print('dist/%s    %d KB　%d 枚' % (os.path.basename(out), os.path.getsize(out) // 1024, len(prs.slides._sldIdLst)))
